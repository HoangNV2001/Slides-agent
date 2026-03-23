"""
PPTX Builder: slide duplication, reordering, deletion, text replacement.
Pure python-pptx + lxml. No external scripts. Runs on any OS.
"""
import copy
import os
import re
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.util import Pt
from lxml import etree


def duplicate_slide(prs: Presentation, source_index: int) -> int:
    """
    Duplicate a slide within a presentation.
    Returns the 0-based index of the new slide (appended at end).

    This works at the lxml/OPC level:
      1. Copy the slide XML element tree
      2. Create a new slide part with a new relationship
      3. Copy all relationships from the source slide (images, layouts, etc.)
      4. Append a new <p:sldId> to the presentation's sldIdLst
    """
    source_slide = prs.slides[source_index]
    slide_layout = source_slide.slide_layout

    # Use the internal API to add a blank slide with the same layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Clear everything from the new slide's spTree
    new_sp_tree = new_slide.shapes._spTree
    for child in list(new_sp_tree):
        tag = etree.QName(child.tag).localname if isinstance(child.tag, str) else ""
        # Keep the cNvGrpSpPr and grpSpPr, remove shapes
        if tag not in ("nvGrpSpPr", "grpSpPr"):
            new_sp_tree.remove(child)

    # Deep-copy all shape elements from source
    source_sp_tree = source_slide.shapes._spTree
    for child in source_sp_tree:
        tag = etree.QName(child.tag).localname if isinstance(child.tag, str) else ""
        if tag not in ("nvGrpSpPr", "grpSpPr"):
            new_sp_tree.append(copy.deepcopy(child))

    # Copy relationships (images, charts, etc.) from source to new slide
    # Build mapping: old rId -> new rId
    rid_map = {}
    for rel in source_slide.part.rels.values():
        # Skip the slide layout relationship (already set)
        if rel.reltype == RT.SLIDE_LAYOUT:
            continue
        try:
            if rel.is_external:
                new_rel = new_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
            else:
                new_rel = new_slide.part.rels.get_or_add(rel.reltype, rel.target_part)
            rid_map[rel.rId] = new_rel if isinstance(new_rel, str) else new_rel
        except Exception:
            # Some relationship types may not copy cleanly; skip gracefully
            pass

    # Update rId references in the new slide XML
    if rid_map:
        _update_rids_in_xml(new_sp_tree, rid_map)

    return len(prs.slides) - 1


def _update_rids_in_xml(element, rid_map: dict):
    """Replace old rId references with new ones in an XML element tree."""
    for attr_name, attr_value in element.attrib.items():
        if attr_value in rid_map:
            element.attrib[attr_name] = rid_map[attr_value]
    for child in element:
        _update_rids_in_xml(child, rid_map)


def delete_slide(prs: Presentation, slide_index: int):
    """
    Delete a slide by index from the presentation.
    Works at the lxml/OPC level.
    """
    slide = prs.slides[slide_index]
    rId = None

    # Find the relationship ID for this slide
    for rel_key, rel in prs.part.rels.items():
        if rel.target_part == slide.part:
            rId = rel.rId
            break

    if rId is None:
        raise ValueError(f"Could not find relationship for slide index {slide_index}")

    # Remove from sldIdLst
    nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    sldIdLst = prs.part._element.find('.//p:sldIdLst', nsmap)
    if sldIdLst is not None:
        for sldId in list(sldIdLst):
            if sldId.get(etree.QName(r_ns, "id")) == rId:
                sldIdLst.remove(sldId)
                break

    # Remove the relationship
    prs.part.drop_rel(rId)


def reorder_slides(prs: Presentation, new_order: List[int]):
    """
    Reorder slides. new_order is a list of current 0-based indices in desired order.
    E.g. [2, 0, 1] means: current slide 2 becomes first, then slide 0, then slide 1.
    """
    nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    sldIdLst = prs.part._element.find('.//p:sldIdLst', nsmap)
    if sldIdLst is None:
        raise ValueError("Could not find sldIdLst in presentation XML")

    sld_ids = list(sldIdLst)

    # Validate
    if sorted(new_order) != list(range(len(sld_ids))):
        raise ValueError(f"new_order must be a permutation of 0..{len(sld_ids)-1}, got {new_order}")

    # Remove all
    for sld_id in sld_ids:
        sldIdLst.remove(sld_id)

    # Re-add in new order
    for idx in new_order:
        sldIdLst.append(sld_ids[idx])


def replace_text_in_slide(slide, replacements: Dict[str, str]):
    """
    Replace text in all shapes of a slide.
    replacements: {old_text: new_text}

    Handles text split across multiple runs by working at the paragraph level.
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            _replace_in_paragraph(paragraph, replacements)


def _shape_top(shape) -> int:
    return int(getattr(shape, "top", 0) or 0)


def _shape_left(shape) -> int:
    return int(getattr(shape, "left", 0) or 0)


def _shape_area(shape) -> int:
    width = int(getattr(shape, "width", 0) or 0)
    height = int(getattr(shape, "height", 0) or 0)
    return width * height


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return _normalize_text_for_match(shape.text_frame.text)


def _is_numeric_marker_text(text: str) -> bool:
    normalized = _normalize_text_for_match(text)
    return bool(re.fullmatch(r"[\dIVXivx]+[.)]?", normalized))


def _shape_is_picture_slot(shape) -> bool:
    shape_type = int(getattr(shape, "shape_type", 0) or 0)
    if shape_type == 13:
        return True

    if getattr(shape, "is_placeholder", False):
        try:
            placeholder_type = str(shape.placeholder_format.type).upper()
            if "PICTURE" in placeholder_type or "CONTENT" in placeholder_type:
                return True
        except Exception:
            return False

    return False


def _find_image_slots(slide) -> List[object]:
    slots = []
    for shape in slide.shapes:
        if _shape_is_picture_slot(shape):
            slots.append(shape)
    return sorted(slots, key=lambda shp: (_shape_top(shp), _shape_left(shp), -_shape_area(shp)))


def _estimate_slide_bounds(slide) -> tuple:
    max_right = 0
    max_bottom = 0
    for shape in slide.shapes:
        max_right = max(max_right, _shape_left(shape) + int(getattr(shape, "width", 0) or 0))
        max_bottom = max(max_bottom, _shape_top(shape) + int(getattr(shape, "height", 0) or 0))
    return max_right or 1, max_bottom or 1


def _looks_like_footer_shape(shape, slide_width: int, slide_height: int) -> bool:
    top = _shape_top(shape)
    bottom = top + int(getattr(shape, "height", 0) or 0)
    width = int(getattr(shape, "width", 0) or 0)
    text = _shape_text(shape)
    if _is_numeric_marker_text(text):
        return True
    if bottom > slide_height * 0.86 and width < slide_width * 0.6:
        return True
    return False


def _dedupe_text_list(values: List[str]) -> List[str]:
    seen = set()
    result = []
    for value in values:
        normalized = _normalize_text_for_match(value)
        if not normalized or normalized in seen:
            continue
        seen.add(normalized)
        result.append(value.strip())
    return result


def _split_body_lines(body: str) -> List[str]:
    if not body:
        return []
    lines = []
    for raw in re.split(r"[\n\r]+", body):
        for piece in raw.split("|"):
            cleaned = piece.strip(" -•\t")
            if cleaned:
                lines.append(cleaned)
    return _dedupe_text_list(lines)


def _build_content_sequence(draft_content: Optional[dict]) -> dict:
    draft_content = draft_content or {}
    bullet_points = draft_content.get("bullet_points") or []
    bullet_points = [str(item).strip() for item in bullet_points if str(item).strip()]
    bullet_points = _dedupe_text_list(bullet_points)

    body_lines = _split_body_lines(draft_content.get("body", ""))
    body_lines = [line for line in body_lines if line not in bullet_points]

    sequence = {
        "title": str(draft_content.get("title") or "").strip(),
        "subtitle": str(draft_content.get("subtitle") or "").strip(),
        "body_lines": bullet_points + body_lines,
    }

    visual_suggestion = str(draft_content.get("visual_suggestion") or "").strip()
    if visual_suggestion and not visual_suggestion.startswith("["):
        sequence["body_lines"].append(visual_suggestion)

    sequence["body_lines"] = _dedupe_text_list(sequence["body_lines"])
    return sequence


def _select_title_shape(shapes: List) -> Optional[object]:
    if not shapes:
        return None

    def score(shape) -> tuple:
        name = (shape.name or "").lower()
        text = _shape_text(shape)
        is_title_placeholder = bool(
            getattr(shape, "is_placeholder", False)
            and "title" in name
        )
        contains_title = "title" in name
        text_len = len(text)
        heading_like = 1 if (8 <= text_len <= 80 and not re.search(r"[,.!?;:|]", text)) else 0
        return (
            1 if is_title_placeholder else 0,
            1 if contains_title else 0,
            heading_like,
            1 if 12 <= text_len <= 120 else 0,
            min(text_len, 80),
            _shape_area(shape),
            -_shape_top(shape),
        )

    return max(shapes, key=score)


def _select_subtitle_shape(shapes: List, title_shape) -> Optional[object]:
    candidates = [shape for shape in shapes if shape != title_shape]
    if not candidates:
        return None

    def score(shape) -> tuple:
        text = _shape_text(shape)
        name = (shape.name or "").lower()
        return (
            1 if "subtitle" in name else 0,
            -_shape_top(shape),
            _shape_area(shape),
            len(text),
        )

    return max(candidates, key=score)


def _set_shape_text(shape, text: str):
    if not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    paragraphs = tf.paragraphs
    lines = [line for line in (text or "").split("\n")] or [""]

    if not paragraphs:
        tf.text = text
        return

    first_para = paragraphs[0]
    _set_paragraph_text(first_para, lines[0])

    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]._p
        p.getparent().remove(p)

    for line in lines[1:]:
        new_para = copy.deepcopy(first_para._p)
        runs = new_para.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}r")
        if runs:
            t_elem = runs[0].find("{http://schemas.openxmlformats.org/drawingml/2006/main}t")
            if t_elem is not None:
                t_elem.text = line
            for r in runs[1:]:
                t_elem = r.find("{http://schemas.openxmlformats.org/drawingml/2006/main}t")
                if t_elem is not None:
                    t_elem.text = ""
        tf._txBody.append(new_para)


def _distribute_lines_across_shapes(lines: List[str], shapes: List[object]) -> List[str]:
    """
    Pack all content lines into the available body shapes instead of dropping overflow.
    Distribution is weighted by approximate shape capacity so small text boxes
    receive less copy than larger body areas.
    Returns one joined text block per shape.
    """
    if not shapes:
        return []
    if not lines:
        return [""] * len(shapes)

    capacities = [max(_estimate_text_capacity(shape), 1) for shape in shapes]
    total_capacity = sum(capacities) or len(shapes)
    remaining_lines = len(lines)
    allocations = []

    for idx, capacity in enumerate(capacities):
        remaining_shapes = len(shapes) - idx
        if idx == len(shapes) - 1:
            take = remaining_lines
        else:
            proportional_take = round(len(lines) * (capacity / total_capacity))
            min_take = 1 if remaining_lines >= remaining_shapes else 0
            max_take = remaining_lines - max(remaining_shapes - 1, 0)
            take = max(min_take, min(proportional_take, max_take))
        allocations.append(take)
        remaining_lines -= take
        total_capacity -= capacity

    blocks = []
    cursor = 0
    for take in allocations:
        chunk = lines[cursor:cursor + take]
        cursor += take
        blocks.append("\n".join(chunk))
    return blocks


def _shape_width_pt(shape) -> float:
    return max(float(int(getattr(shape, "width", 0) or 0)) / 12700.0, 1.0)


def _shape_height_pt(shape) -> float:
    return max(float(int(getattr(shape, "height", 0) or 0)) / 12700.0, 1.0)


def _iter_shape_runs(shape):
    if not getattr(shape, "has_text_frame", False):
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            yield run


def _get_shape_font_sizes(shape) -> List[float]:
    sizes = []
    for run in _iter_shape_runs(shape):
        if run.font.size:
            sizes.append(float(run.font.size.pt))
    return sizes


def _get_base_font_size(shape) -> float:
    sizes = _get_shape_font_sizes(shape)
    if sizes:
        return max(min(sum(sizes) / len(sizes), 28.0), 10.0)
    return 18.0


def _estimate_text_capacity(shape) -> int:
    width_pt = _shape_width_pt(shape)
    height_pt = _shape_height_pt(shape)
    font_size = _get_base_font_size(shape)
    line_height = max(font_size * 1.25, 12.0)
    chars_per_line = max(int(width_pt / max(font_size * 0.58, 1.0)), 6)
    line_capacity = max(int(height_pt / line_height), 1)
    return max(chars_per_line * line_capacity, 1)


def _estimate_overflow_ratio(shape) -> float:
    text = _shape_text(shape)
    if not text:
        return 0.0
    capacity = _estimate_text_capacity(shape)
    weighted_length = max(len(text), 0) + max(text.count("\n"), 0) * 18
    return weighted_length / max(capacity, 1)


def _set_shape_autofit(shape):
    if not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    try:
        tf.vertical_anchor = MSO_ANCHOR.TOP
    except Exception:
        pass


def _scale_shape_font_sizes(shape, scale: float, min_size: float = 10.0):
    scale = max(min(scale, 1.0), 0.45)
    base_size = _get_base_font_size(shape)
    updated = 0
    for run in _iter_shape_runs(shape):
        current_size = float(run.font.size.pt) if run.font.size else base_size
        target_size = max(min_size, round(current_size * scale, 1))
        if abs(target_size - current_size) < 0.2:
            continue
        run.font.size = Pt(target_size)
        updated += 1
    return updated


def _remove_shape(shape):
    element = getattr(shape, "_element", None)
    if element is None:
        return False
    parent = element.getparent()
    if parent is None:
        return False
    parent.remove(element)
    return True


def _is_removable_text_shape(shape, slide_width: int, slide_height: int) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False
    text = _shape_text(shape)
    if text:
        return False
    if _looks_like_footer_shape(shape, slide_width, slide_height):
        return False
    name = (getattr(shape, "name", "") or "").lower()
    if any(token in name for token in ("date", "footer", "slide number")):
        return False
    if getattr(shape, "is_placeholder", False):
        return True
    return _shape_area(shape) < (slide_width * slide_height * 0.08)


def review_and_fix_slide(slide, slide_number: int) -> dict:
    """
    Inspect a populated slide and apply lightweight layout hygiene fixes.
    This aims to keep slides presentation-ready without changing slide meaning.
    """
    slide_width, slide_height = _estimate_slide_bounds(slide)
    actions = []
    issues = []

    removable_shapes = []
    for shape in list(slide.shapes):
        if not getattr(shape, "has_text_frame", False):
            continue
        text = _shape_text(shape)
        if text:
            _set_shape_autofit(shape)
            overflow_ratio = _estimate_overflow_ratio(shape)
            if overflow_ratio > 1.15:
                scale = min(1.0 / overflow_ratio, 0.9)
                updated_runs = _scale_shape_font_sizes(shape, scale)
                if updated_runs:
                    actions.append(
                        f"shrunk text in '{shape.name}' to reduce overflow risk ({overflow_ratio:.2f}x)"
                    )
                if _estimate_overflow_ratio(shape) > 1.2:
                    issues.append(
                        f"text in '{shape.name}' may still be dense after autofit"
                    )
        elif _is_removable_text_shape(shape, slide_width, slide_height):
            removable_shapes.append(shape)

    for shape in removable_shapes:
        if _remove_shape(shape):
            actions.append(f"removed unused placeholder '{shape.name}'")

    return {
        "slide_number": slide_number,
        "actions": actions,
        "issues": issues,
    }


def _choose_slide_images(draft_content: Optional[dict], document_images: Optional[List[dict]]) -> List[dict]:
    draft_content = draft_content or {}
    document_images = document_images or []
    image_by_id = {
        str(image.get("id")): image for image in document_images
        if image.get("id") and image.get("path")
    }

    selected = []
    for image_id in draft_content.get("source_image_ids") or []:
        image = image_by_id.get(str(image_id))
        if image and os.path.exists(image.get("path", "")):
            selected.append(image)
    if selected:
        return selected

    visual_suggestion = str(draft_content.get("visual_suggestion") or "")
    if "[IMAGE" not in visual_suggestion.upper():
        return []

    best_match = _find_best_matching_image(draft_content, document_images)
    return [best_match] if best_match else []

    return selected


def _find_best_matching_image(draft_content: Optional[dict], document_images: Optional[List[dict]]) -> Optional[dict]:
    document_images = document_images or []
    slide_text = " ".join(
        filter(
            None,
            [
                str((draft_content or {}).get("title") or ""),
                str((draft_content or {}).get("subtitle") or ""),
                str((draft_content or {}).get("body") or ""),
                " ".join(str(x) for x in ((draft_content or {}).get("bullet_points") or [])),
                str((draft_content or {}).get("visual_suggestion") or ""),
            ],
        )
    ).lower()
    slide_keywords = set(_extract_keywords(slide_text))
    if not slide_keywords:
        return None

    best_image = None
    best_score = 0
    for image in document_images:
        image_path = image.get("path", "")
        if not image_path or not os.path.exists(image_path):
            continue
        image_keywords = set(image.get("context_keywords") or _extract_keywords(
            " ".join(filter(None, [image.get("caption", ""), image.get("nearby_text", "")]))
        ))
        overlap = slide_keywords & image_keywords
        score = len(overlap)
        if score > best_score:
            best_score = score
            best_image = image

    return best_image if best_score >= 2 else None


def _apply_images_to_slide(slide, draft_content: Optional[dict], document_images: Optional[List[dict]]) -> int:
    """Overlay extracted source images onto image slots in the template slide."""
    selected_images = _choose_slide_images(draft_content, document_images)
    if not selected_images:
        return 0

    image_slots = _find_image_slots(slide)
    if not image_slots:
        return 0

    inserted = 0
    for slot, image in zip(image_slots, selected_images):
        image_path = image.get("path", "")
        if not image_path or not os.path.exists(image_path):
            continue

        try:
            if getattr(slot, "has_text_frame", False):
                _clear_shape_text(slot)
            slide.shapes.add_picture(
                image_path,
                slot.left,
                slot.top,
                width=slot.width,
                height=slot.height,
            )
            inserted += 1
        except Exception:
            continue

    return inserted


def _apply_draft_content_to_slide(slide, draft_content: Optional[dict]) -> int:
    """
    Dynamically assign drafted content into the text shapes available on a slide.
    Unused text boxes are cleared so template placeholders do not remain visible.
    """
    content = _build_content_sequence(draft_content)
    if not any([content["title"], content["subtitle"], content["body_lines"]]):
        return 0

    text_shapes = [
        shape for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and _shape_text(shape)
    ]
    if not text_shapes:
        return 0

    slide_width, slide_height = _estimate_slide_bounds(slide)

    updates = 0
    used_shape_ids = set()

    title_shape = _select_title_shape(text_shapes)
    if title_shape and content["title"]:
        _set_shape_text(title_shape, content["title"])
        used_shape_ids.add(title_shape.shape_id)
        updates += 1

    subtitle_shape = None
    if content["subtitle"]:
        subtitle_shape = _select_subtitle_shape(text_shapes, title_shape)
        if subtitle_shape:
            _set_shape_text(subtitle_shape, content["subtitle"])
            used_shape_ids.add(subtitle_shape.shape_id)
            updates += 1

    remaining_shapes = [
        shape for shape in sorted(text_shapes, key=lambda shp: (_shape_top(shp), _shape_left(shp), -_shape_area(shp)))
        if shape.shape_id not in used_shape_ids
    ]

    body_shapes = []
    for shape in remaining_shapes:
        shape_text = _shape_text(shape)
        if not shape_text or _looks_like_footer_shape(shape, slide_width, slide_height):
            continue
        body_shapes.append(shape)

    body_lines = list(content["body_lines"])
    body_blocks = _distribute_lines_across_shapes(body_lines, body_shapes)
    for idx, shape in enumerate(body_shapes):
        block_text = body_blocks[idx] if idx < len(body_blocks) else ""
        if block_text:
            _set_shape_text(shape, block_text)
        else:
            _clear_shape_text(shape)
        used_shape_ids.add(shape.shape_id)
        updates += 1

    for shape in remaining_shapes:
        if shape.shape_id in used_shape_ids:
            continue
        shape_text = _shape_text(shape)
        if shape_text and not _is_numeric_marker_text(shape_text) and _looks_like_placeholder_text(shape_text):
            _clear_shape_text(shape)
            updates += 1

    return updates


PLACEHOLDER_PATTERNS = [
    re.compile(r"\blorem ipsum\b", re.IGNORECASE),
    re.compile(r"\bconsectetur adipiscing\b", re.IGNORECASE),
    re.compile(r"\bsed do eiusmod\b", re.IGNORECASE),
    re.compile(r"\but enim ad minim\b", re.IGNORECASE),
    re.compile(r"\bquis nostrud exercitation\b", re.IGNORECASE),
    re.compile(r"\bduis aute irure\b", re.IGNORECASE),
    re.compile(r"\bexcepteur sint occaecat\b", re.IGNORECASE),
    re.compile(r"\bclick to add\b", re.IGNORECASE),
    re.compile(r"\bdouble click to edit\b", re.IGNORECASE),
    re.compile(r"\byour (title|subtitle|text|content|company|name)\b", re.IGNORECASE),
    re.compile(r"\b(insert|add) (text|title|subtitle|content|description|agenda|date|name)\b", re.IGNORECASE),
    re.compile(r"\b(sample|dummy|placeholder) (text|content|copy)\b", re.IGNORECASE),
    re.compile(r"\bagenda item\b", re.IGNORECASE),
    re.compile(r"\bsubheading\b", re.IGNORECASE),
    re.compile(r"\btext here\b", re.IGNORECASE),
    re.compile(r"\btitle here\b", re.IGNORECASE),
    re.compile(r"\bsubtitle here\b", re.IGNORECASE),
    re.compile(r"\bdescription here\b", re.IGNORECASE),
    re.compile(r"\bxxx+\b", re.IGNORECASE),
    re.compile(r"\b(todo|tbd)\b", re.IGNORECASE),
    re.compile(r"\b(đảm bảo|bao dam) hệ font thương hiệu\b", re.IGNORECASE),
    re.compile(r"\bsử dụng format chuẩn của template\b", re.IGNORECASE),
    re.compile(r"\bquy chuẩn font chữ tiêu đề\b", re.IGNORECASE),
    re.compile(r"\bsử dụng đúng bộ màu quy định\b", re.IGNORECASE),
    re.compile(r"\bkhông sử dụng thêm bất cứ font hỗ trợ nào khác\b", re.IGNORECASE),
    re.compile(r"\bnên sử dụng nhiều hình ảnh để hỗ trợ cho nội dung thông tin\b", re.IGNORECASE),
    re.compile(r"\bchuyển sang định dạng pdf khi hoàn thành\b", re.IGNORECASE),
    re.compile(r"\bsử dụng đúng màu sắc thương hiệu\b", re.IGNORECASE),
]


PLACEHOLDER_SUBSTITUTIONS = [
    re.compile(
        r",?\s*lorem ipsum.*$",
        re.IGNORECASE,
    ),
    re.compile(
        r",?\s*consectetur adipiscing elit.*$",
        re.IGNORECASE,
    ),
    re.compile(
        r"\b(lorem ipsum|consectetur adipiscing|sed do eiusmod|ut enim ad minim|quis nostrud exercitation|duis aute irure|excepteur sint occaecat).*",
        re.IGNORECASE,
    ),
]


def _normalize_text_for_match(text: str) -> str:
    return re.sub(r"\s+", " ", (text or "").strip())


def _looks_like_placeholder_text(text: str) -> bool:
    normalized = _normalize_text_for_match(text)
    if not normalized:
        return False
    return any(pattern.search(normalized) for pattern in PLACEHOLDER_PATTERNS)


def _strip_placeholder_substrings(text: str) -> str:
    cleaned = text or ""
    for pattern in PLACEHOLDER_SUBSTITUTIONS:
        cleaned = pattern.sub("", cleaned)
    cleaned = re.sub(r"\s+", " ", cleaned).strip(" ,;:-")
    return cleaned


def _should_clear_entire_text(text: str) -> bool:
    normalized = _normalize_text_for_match(text)
    if not normalized:
        return False

    color_code_count = len(re.findall(r"#[0-9A-Fa-f]{6}", normalized))
    rgb_count = len(re.findall(r"\bRGB\s*:", normalized, re.IGNORECASE))

    if color_code_count >= 2 or rgb_count >= 2:
        return True

    return False


def _clear_paragraph(paragraph):
    for run in paragraph.runs:
        run.text = ""


def _clear_shape_text(shape):
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        _clear_paragraph(paragraph)


def _set_paragraph_text(paragraph, text: str):
    runs = paragraph.runs
    if runs:
        runs[0].text = text
        for run in runs[1:]:
            run.text = ""
    else:
        paragraph.text = text


def clear_unused_placeholder_text(slide, replacements: Optional[Dict[str, str]] = None) -> int:
    """
    Remove obvious template placeholder text left on a slide after replacements.
    Only clears text that still looks like generic placeholder copy.
    """
    cleared = 0
    applied_replacement_keys = {
        _normalize_text_for_match(key) for key in (replacements or {}) if key
    }

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        shape_text = _normalize_text_for_match(shape.text_frame.text)
        if shape_text and shape_text in applied_replacement_keys:
            continue
        if _should_clear_entire_text(shape_text):
            _clear_shape_text(shape)
            cleared += 1
            continue

        for paragraph in shape.text_frame.paragraphs:
            paragraph_text = _normalize_text_for_match("".join(run.text for run in paragraph.runs))
            if not paragraph_text:
                continue
            if paragraph_text in applied_replacement_keys:
                continue
            if _should_clear_entire_text(paragraph_text):
                _clear_paragraph(paragraph)
                cleared += 1
                continue
            stripped_text = _strip_placeholder_substrings(paragraph_text)
            if stripped_text != paragraph_text:
                if stripped_text:
                    _set_paragraph_text(paragraph, stripped_text)
                else:
                    _clear_paragraph(paragraph)
                cleared += 1
                continue
            if _looks_like_placeholder_text(paragraph_text):
                _clear_paragraph(paragraph)
                cleared += 1

    return cleared


def _replace_in_paragraph(paragraph, replacements: Dict[str, str]):
    """
    Replace text in a paragraph, handling text split across runs.

    Strategy:
    1. Concatenate all run texts to get full paragraph text.
    2. If a replacement matches, redistribute the new text across runs,
       preserving the formatting of the first matching run.
    """
    full_text = "".join(run.text for run in paragraph.runs)
    if not full_text:
        return

    new_full_text = full_text
    changed = False
    for old_text, new_text in replacements.items():
        if old_text in new_full_text:
            new_full_text = new_full_text.replace(old_text, new_text)
            changed = True

    if not changed:
        return

    # Redistribute: put all text in first run, clear others
    runs = paragraph.runs
    if not runs:
        return

    runs[0].text = new_full_text
    for run in runs[1:]:
        run.text = ""


def replace_all_text_by_shape_name(slide, shape_name: str, new_text: str):
    """Replace ALL text in a shape identified by name."""
    for shape in slide.shapes:
        if shape.name == shape_name and shape.has_text_frame:
            tf = shape.text_frame
            # Keep formatting from first paragraph/run
            if tf.paragraphs:
                first_para = tf.paragraphs[0]
                # Put new text lines into paragraphs
                lines = new_text.split("\n")

                # Set first paragraph
                if first_para.runs:
                    first_para.runs[0].text = lines[0]
                    for r in first_para.runs[1:]:
                        r.text = ""
                else:
                    first_para.text = lines[0]

                # Add remaining lines as new paragraphs
                for line in lines[1:]:
                    new_para = copy.deepcopy(first_para._p)
                    # Clear runs in copied paragraph and set text
                    runs = new_para.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}r")
                    if runs:
                        t_elem = runs[0].find("{http://schemas.openxmlformats.org/drawingml/2006/main}t")
                        if t_elem is not None:
                            t_elem.text = line
                        for r in runs[1:]:
                            t_elem = r.find("{http://schemas.openxmlformats.org/drawingml/2006/main}t")
                            if t_elem is not None:
                                t_elem.text = ""
                    tf._txBody.append(new_para)

            return True
    return False


def get_slide_text_inventory(slide) -> List[dict]:
    """Get all text shapes and their content from a slide."""
    inventory = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                info = {
                    "shape_name": shape.name,
                    "shape_id": shape.shape_id,
                    "text": text,
                    "is_placeholder": shape.is_placeholder,
                }
                if shape.is_placeholder:
                    info["placeholder_idx"] = shape.placeholder_format.idx
                inventory.append(info)
    return inventory


def build_presentation_from_plan(
    template_path: str,
    slide_plan: List[dict],
    output_path: str,
) -> dict:
    """
    High-level: build a full presentation from a template + slide plan.

    slide_plan items:
    {
        "source_slide_index": int,     # 0-based index in template
        "text_replacements": {str: str},  # old_text -> new_text
        "draft_content": { ... },         # slide content from drafter
        "document_images": [ ... ],       # extracted source images
    }

    Returns dict with status, steps, warnings.
    """
    result = {"status": "starting", "steps": [], "warnings": [], "review_report": []}

    try:
        prs = Presentation(template_path)
        original_count = len(prs.slides)
        result["steps"].append(f"Loaded template: {original_count} slides")

        # Phase 1: Duplicate slides as needed.
        # We need len(slide_plan) slides total.
        # Track: plan_index -> actual slide index in the growing presentation.
        slide_index_map = {}  # plan_item_index -> slide_index in prs

        # Count usage of each source slide
        source_usage = {}
        for i, item in enumerate(slide_plan):
            src = item["source_slide_index"]
            if src not in source_usage:
                source_usage[src] = []
            source_usage[src].append(i)

        for src_idx, plan_indices in source_usage.items():
            if src_idx >= original_count:
                result["warnings"].append(f"Source slide index {src_idx} out of range (template has {original_count})")
                continue

            # First usage: use the original slide
            slide_index_map[plan_indices[0]] = src_idx

            # Additional usages: duplicate
            for extra_plan_idx in plan_indices[1:]:
                try:
                    new_idx = duplicate_slide(prs, src_idx)
                    slide_index_map[extra_plan_idx] = new_idx
                    result["steps"].append(f"  Duplicated slide {src_idx} → {new_idx}")
                except Exception as e:
                    result["warnings"].append(f"Failed to duplicate slide {src_idx}: {e}")
                    slide_index_map[extra_plan_idx] = src_idx  # fallback

        result["steps"].append(f"✓ Slide structure ready ({len(prs.slides)} total slides)")

        # Phase 2: Apply text replacements
        result["steps"].append("Applying text replacements...")
        for plan_idx, item in enumerate(slide_plan):
            actual_idx = slide_index_map.get(plan_idx)
            if actual_idx is None:
                result["warnings"].append(f"No slide mapped for plan item {plan_idx}")
                continue

            try:
                slide = prs.slides[actual_idx]
                replacements = item.get("text_replacements", {})
                if replacements:
                    replace_text_in_slide(slide, replacements)
                dynamic_updates = _apply_draft_content_to_slide(slide, item.get("draft_content", {}))
                inserted_images = _apply_images_to_slide(
                    slide,
                    item.get("draft_content", {}),
                    item.get("document_images", []),
                )
                cleared = clear_unused_placeholder_text(slide, replacements)
                review = review_and_fix_slide(slide, plan_idx + 1)
                result["review_report"].append(review)
                step = f"  Slide {actual_idx}: {len(replacements)} replacements applied"
                if dynamic_updates:
                    step += f", {dynamic_updates} dynamic text slot(s) updated"
                if inserted_images:
                    step += f", {inserted_images} image(s) inserted"
                if cleared:
                    step += f", {cleared} placeholder text block(s) cleared"
                if review["actions"]:
                    step += f", review fixed {len(review['actions'])} item(s)"
                result["steps"].append(step)
                for issue in review["issues"]:
                    result["warnings"].append(f"Slide {plan_idx + 1}: {issue}")
            except Exception as e:
                result["warnings"].append(f"Error replacing text in slide {actual_idx}: {e}")

        result["steps"].append("✓ Text replacements applied")

        # Phase 3: Reorder slides to match plan order, and remove unused slides
        result["steps"].append("Reordering slides...")
        desired_order = [slide_index_map[i] for i in range(len(slide_plan)) if i in slide_index_map]

        # We need to keep only the slides in desired_order.
        # To avoid index shifting issues, we reorder first then delete extras.
        all_indices = set(range(len(prs.slides)))
        used_indices = set(desired_order)
        unused_indices = all_indices - used_indices

        # Build full reorder: desired slides first, then unused
        full_order = desired_order + sorted(unused_indices)
        try:
            reorder_slides(prs, full_order)
            result["steps"].append(f"✓ Slides reordered")
        except Exception as e:
            result["warnings"].append(f"Reorder warning: {e}")

        # Delete unused slides (they're now at the end)
        # Delete from the back to avoid index shifting
        num_to_keep = len(desired_order)
        num_total = len(prs.slides)
        for i in range(num_total - 1, num_to_keep - 1, -1):
            try:
                delete_slide(prs, i)
            except Exception as e:
                result["warnings"].append(f"Could not remove unused slide {i}: {e}")

        result["steps"].append(f"✓ Final presentation: {len(prs.slides)} slides")

        # Phase 4: Save
        prs.save(output_path)
        result["steps"].append(f"✓ Saved to {os.path.basename(output_path)}")

        # Phase 5: Validate
        result["steps"].append("Validating output...")
        try:
            val_prs = Presentation(output_path)
            val_text = ""
            for i, slide in enumerate(val_prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        slide_text.append(shape.text_frame.text)
                val_text += f"--- Slide {i + 1} ---\n" + "\n".join(slide_text) + "\n\n"

            # Check for leftover placeholders
            placeholders = re.findall(
                r"\b[Xx]{3,}\b|lorem|ipsum|\bTODO\b|\bTBD\b|\[insert|click to add|placeholder|sample text|text here",
                val_text, re.IGNORECASE
            )
            if placeholders:
                result["warnings"].append(f"Possible leftover placeholders: {placeholders[:5]}")
            result["validation_text"] = val_text[:3000]
        except Exception as e:
            result["warnings"].append(f"Validation error: {e}")

        result["status"] = "success"
        result["steps"].append("✓ Generation complete!")

    except Exception as e:
        result["status"] = "error"
        result["error"] = str(e)
        import traceback
        result["traceback"] = traceback.format_exc()

    return result
