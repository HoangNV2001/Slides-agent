"""
Template analyzer: extracts structure, text inventory, and slide metadata from PPTX templates.
Uses python-pptx only — no external scripts required.
"""
import os
from collections import Counter
from typing import Optional
from pptx import Presentation
from pptx.util import Inches, Pt, Emu


def analyze_template(template_path: str) -> dict:
    """
    Full template analysis using python-pptx.
    Returns dict with slides list, total count, text inventory, layout info.
    """
    result = {
        "template_path": template_path,
        "slides": [],
        "total_slides": 0,
        "text_inventory": {},
        "slide_layouts_available": [],
        "mapping_slides": [],
    }

    try:
        prs = Presentation(template_path)

        # Collect available layouts
        for layout in prs.slide_layouts:
            result["slide_layouts_available"].append({
                "name": layout.name,
                "placeholders": [
                    {"idx": ph.placeholder_format.idx, "name": ph.name, "type": str(ph.placeholder_format.type)}
                    for ph in layout.placeholders
                ],
            })

        slide_width = int(prs.slide_width)
        slide_height = int(prs.slide_height)
        result["slide_width"] = slide_width
        result["slide_height"] = slide_height

        all_texts = []

        # Analyze each slide
        for slide_idx, slide in enumerate(prs.slides):
            slide_info = {
                "index": slide_idx,
                "slide_id": slide.slide_id,
                "layout_name": slide.slide_layout.name if slide.slide_layout else "Unknown",
                "text_shapes": [],
                "image_shapes": [],
                "has_images": False,
                "has_charts": False,
                "has_tables": False,
                "shape_count": len(slide.shapes),
            }

            for shape in slide.shapes:
                # Check shape types
                if shape.shape_type is not None:
                    type_val = int(shape.shape_type)
                    if type_val == 13:  # Picture
                        slide_info["has_images"] = True
                        slide_info["image_shapes"].append({
                            "shape_name": shape.name,
                            "shape_id": shape.shape_id,
                            "left": int(shape.left),
                            "top": int(shape.top),
                            "width": int(shape.width),
                            "height": int(shape.height),
                        })
                    elif type_val == 3:  # Chart
                        slide_info["has_charts"] = True

                if shape.has_table:
                    slide_info["has_tables"] = True

                if shape.has_text_frame:
                    text_content = []
                    for para in shape.text_frame.paragraphs:
                        para_text = para.text.strip()
                        if para_text:
                            text_content.append(para_text)

                    if text_content:
                        shape_data = {
                            "shape_name": shape.name,
                            "shape_id": shape.shape_id,
                            "text": "\n".join(text_content),
                            "is_placeholder": shape.is_placeholder,
                            "left": int(shape.left),
                            "top": int(shape.top),
                            "width": int(shape.width),
                            "height": int(shape.height),
                        }
                        if shape.is_placeholder:
                            shape_data["placeholder_idx"] = shape.placeholder_format.idx
                            shape_data["placeholder_type"] = str(shape.placeholder_format.type)
                        slide_info["text_shapes"].append(shape_data)
                        all_texts.append(_normalize_text(shape_data["text"]))

            result["slides"].append(slide_info)
            result["text_inventory"][f"slide_{slide_idx}"] = slide_info["text_shapes"]

        result["total_slides"] = len(result["slides"])
        repeated_texts = {
            text for text, count in Counter(t for t in all_texts if t).items()
            if count >= max(2, int(result["total_slides"] * 0.4))
        }

        for slide in result["slides"]:
            _annotate_slide_roles(slide, slide_width, slide_height, repeated_texts)

        result["mapping_slides"] = _build_mapping_slide_inventory(result["slides"])

    except Exception as e:
        result["error"] = str(e)

    return result


def get_template_summary(analysis: dict) -> str:
    """Generate a human-readable summary of the template."""
    lines = [
        f"Template: {os.path.basename(analysis['template_path'])}",
        f"Total slides: {analysis['total_slides']}",
    ]

    # Available layouts
    layouts = analysis.get("slide_layouts_available", [])
    if layouts:
        lines.append(f"Available layouts: {', '.join(l['name'] for l in layouts)}")

    lines.append("")
    lines.append("Slide inventory:")

    slides_for_mapping = analysis.get("mapping_slides") or analysis.get("slides", [])
    for slide in slides_for_mapping:
        idx = slide["index"]
        layout = slide.get("layout_name", "?")
        simplified = slide.get("simplified_layout", {})
        extras = []
        if slide.get("has_images"):
            extras.append("has images")
        if slide.get("has_charts"):
            extras.append("has charts")
        if slide.get("has_tables"):
            extras.append("has tables")
        extra_str = f" ({', '.join(extras)})" if extras else ""

        lines.append(f"\n  Slide {idx} [layout: {layout}]{extra_str}:")
        lines.append(
            "    - "
            + f"type={simplified.get('layout_kind', 'unknown')}, "
            + f"title_slots={simplified.get('title_slots', 0)}, "
            + f"body_slots={simplified.get('body_slots', 0)}, "
            + f"image_slots={simplified.get('image_slots', 0)}"
        )
        visual = slide.get("visual_layout", {})
        if visual:
            lines.append(
                "    - "
                + f"visual_family={visual.get('layout_family', 'unknown')}, "
                + f"columns={visual.get('column_count', 0)}, "
                + f"image_alignment={visual.get('image_alignment', 'none')}, "
                + f"emphasis={visual.get('emphasis', 'balanced')}"
            )
        for slot in simplified.get("slots", [])[:6]:
            lines.append(
                f"    - [{slot.get('role')}] {slot.get('shape_name')}: "
                f"\"{slot.get('text_preview', '')}\""
            )

    return "\n".join(lines)


def _normalize_text(text: str) -> str:
    return " ".join((text or "").split()).strip().lower()


def _annotate_slide_roles(slide: dict, slide_width: int, slide_height: int, repeated_texts: set):
    slots = []
    for shape in slide.get("text_shapes", []):
        role = _classify_text_role(shape, slide_width, slide_height, repeated_texts)
        shape["role"] = role
        if role not in {"header", "footer", "decorative"}:
            slots.append({
                "shape_id": shape["shape_id"],
                "shape_name": shape["shape_name"],
                "role": role,
                "text_preview": shape["text"][:80].replace("\n", " | "),
                "left": shape.get("left", 0),
                "top": shape.get("top", 0),
                "width": shape.get("width", 0),
                "height": shape.get("height", 0),
            })

    image_slots = slide.get("image_shapes", [])
    for img in image_slots:
        slots.append({
            "shape_id": img["shape_id"],
            "shape_name": img["shape_name"],
            "role": "image",
            "text_preview": "",
            "left": img.get("left", 0),
            "top": img.get("top", 0),
            "width": img.get("width", 0),
            "height": img.get("height", 0),
        })

    title_slots = len([slot for slot in slots if slot["role"] == "title"])
    body_slots = len([slot for slot in slots if slot["role"] in {"subtitle", "body", "quote", "comparison"}])
    image_slot_count = len([slot for slot in slots if slot["role"] == "image"])

    if image_slot_count and body_slots >= 1:
        layout_kind = "image_text"
    elif body_slots >= 2:
        layout_kind = "multi_text"
    elif title_slots and body_slots:
        layout_kind = "title_text"
    elif body_slots == 1:
        layout_kind = "single_text"
    else:
        layout_kind = "title_only"

    slide["simplified_layout"] = {
        "layout_kind": layout_kind,
        "title_slots": title_slots,
        "body_slots": body_slots,
        "image_slots": image_slot_count,
        "slots": slots,
    }
    slide["visual_layout"] = _build_visual_layout(slots, slide_width, slide_height, layout_kind)


def _classify_text_role(shape: dict, slide_width: int, slide_height: int, repeated_texts: set) -> str:
    top = shape.get("top", 0)
    left = shape.get("left", 0)
    width = shape.get("width", 0)
    height = shape.get("height", 0)
    text = shape.get("text", "")
    normalized = _normalize_text(text)
    name = (shape.get("shape_name") or "").lower()
    bottom = top + height

    if normalized in repeated_texts and bottom > slide_height * 0.82:
        return "footer"
    if normalized in repeated_texts and top < slide_height * 0.12:
        return "header"
    if bottom > slide_height * 0.88:
        return "footer"
    if top < slide_height * 0.08 and width < slide_width * 0.35:
        return "header"
    if len(normalized) <= 3:
        return "decorative"
    if shape.get("is_placeholder") and "title" in name:
        return "title"
    if top < slide_height * 0.28 and width > slide_width * 0.35 and len(normalized) <= 120:
        return "title"
    if "quote" in name or len(normalized) > 140:
        return "quote"
    if width < slide_width * 0.42:
        return "comparison"
    if top < slide_height * 0.45 and len(normalized) <= 120:
        return "subtitle"
    return "body"


def _bucket_band(value: float) -> str:
    if value < 0.22:
        return "top"
    if value < 0.58:
        return "middle"
    return "bottom"


def _bucket_zone(start: float, end: float) -> str:
    if start <= 0.08 and end >= 0.92:
        return "full"
    center = (start + end) / 2.0
    if center < 0.34:
        return "left"
    if center > 0.66:
        return "right"
    return "center"


def _cluster_count(values: list, threshold: float = 0.18) -> int:
    ordered = sorted(v for v in values if v is not None)
    if not ordered:
        return 0
    groups = [ordered[0]]
    count = 1
    for value in ordered[1:]:
        if abs(value - groups[-1]) > threshold:
            count += 1
            groups.append(value)
    return count


def _describe_slot_geometry(slot: dict, slide_width: int, slide_height: int) -> dict:
    left = int(slot.get("left", 0) or 0)
    top = int(slot.get("top", 0) or 0)
    width = max(int(slot.get("width", 0) or 0), 1)
    height = max(int(slot.get("height", 0) or 0), 1)
    x0 = left / max(slide_width, 1)
    y0 = top / max(slide_height, 1)
    x1 = (left + width) / max(slide_width, 1)
    y1 = (top + height) / max(slide_height, 1)
    return {
        "x": round(x0, 3),
        "y": round(y0, 3),
        "w": round(width / max(slide_width, 1), 3),
        "h": round(height / max(slide_height, 1), 3),
        "vertical_band": _bucket_band(y0),
        "horizontal_zone": _bucket_zone(x0, x1),
        "area_ratio": round((width * height) / max(slide_width * slide_height, 1), 3),
    }


def _infer_layout_family(layout_kind: str, title_slots: list, body_slots: list, image_slots: list) -> str:
    title_is_hero = any(slot["geometry"]["area_ratio"] >= 0.14 for slot in title_slots)
    body_columns = _cluster_count([slot["geometry"]["x"] for slot in body_slots])
    image_columns = _cluster_count([slot["geometry"]["x"] for slot in image_slots])

    if image_slots and body_slots:
        if any(slot["geometry"]["horizontal_zone"] == "full" for slot in image_slots):
            return "full_bleed_image_text"
        if image_columns >= 2 and body_columns >= 2:
            return "grid_mix"
        if any(slot["geometry"]["horizontal_zone"] == "left" for slot in image_slots) and any(
            slot["geometry"]["horizontal_zone"] in {"center", "right"} for slot in body_slots
        ):
            return "image_left_text_right"
        if any(slot["geometry"]["horizontal_zone"] == "right" for slot in image_slots) and any(
            slot["geometry"]["horizontal_zone"] in {"left", "center"} for slot in body_slots
        ):
            return "text_left_image_right"
        return "stacked_image_text"

    if body_columns >= 2:
        return "two_column_text"
    if title_is_hero and not body_slots:
        return "hero_title"
    if title_slots and body_slots:
        return "title_body"
    if body_slots:
        return "body_only"
    return layout_kind or "unknown"


def _build_visual_layout(slots: list, slide_width: int, slide_height: int, layout_kind: str) -> dict:
    enriched_slots = []
    title_slots = []
    body_slots = []
    image_slots = []

    for slot in slots:
        geometry = _describe_slot_geometry(slot, slide_width, slide_height)
        enriched = {
            "shape_id": slot.get("shape_id"),
            "shape_name": slot.get("shape_name"),
            "role": slot.get("role"),
            "text_preview": slot.get("text_preview", ""),
            "geometry": geometry,
        }
        enriched_slots.append(enriched)
        if slot.get("role") == "title":
            title_slots.append(enriched)
        elif slot.get("role") == "image":
            image_slots.append(enriched)
        else:
            body_slots.append(enriched)

    visual_layout = {
        "layout_family": _infer_layout_family(layout_kind, title_slots, body_slots, image_slots),
        "column_count": max(_cluster_count([slot["geometry"]["x"] for slot in body_slots]), 1 if body_slots else 0),
        "image_alignment": "none",
        "emphasis": "balanced",
        "slot_map": enriched_slots,
        "reading_order": [
            {
                "role": slot["role"],
                "shape_name": slot["shape_name"],
                "band": slot["geometry"]["vertical_band"],
                "zone": slot["geometry"]["horizontal_zone"],
            }
            for slot in sorted(
                enriched_slots,
                key=lambda item: (
                    item["geometry"]["y"],
                    item["geometry"]["x"],
                    -item["geometry"]["area_ratio"],
                ),
            )
        ],
    }

    if image_slots:
        if any(slot["geometry"]["horizontal_zone"] == "left" for slot in image_slots):
            visual_layout["image_alignment"] = "left"
        elif any(slot["geometry"]["horizontal_zone"] == "right" for slot in image_slots):
            visual_layout["image_alignment"] = "right"
        elif any(slot["geometry"]["horizontal_zone"] == "full" for slot in image_slots):
            visual_layout["image_alignment"] = "full_bleed"
        else:
            visual_layout["image_alignment"] = "center"

    title_area = sum(slot["geometry"]["area_ratio"] for slot in title_slots)
    body_area = sum(slot["geometry"]["area_ratio"] for slot in body_slots)
    image_area = sum(slot["geometry"]["area_ratio"] for slot in image_slots)
    if title_area >= max(body_area, image_area) * 1.1 and title_area >= 0.12:
        visual_layout["emphasis"] = "title"
    elif image_area >= max(title_area, body_area) * 1.1 and image_area >= 0.16:
        visual_layout["emphasis"] = "visual"
    elif body_area >= max(title_area, image_area) * 1.1:
        visual_layout["emphasis"] = "content"

    return visual_layout


def _build_mapping_slide_inventory(slides: list) -> list:
    selected = []
    seen_fingerprints = set()
    for slide in slides:
        simplified = slide.get("simplified_layout", {})
        visual = slide.get("visual_layout", {})
        role_signature = tuple(
            (
                slot.get("role"),
                slot.get("geometry", {}).get("vertical_band"),
                slot.get("geometry", {}).get("horizontal_zone"),
            )
            for slot in visual.get("slot_map", [])
        )
        fingerprint = (
            simplified.get("layout_kind"),
            simplified.get("title_slots", 0),
            simplified.get("body_slots", 0),
            simplified.get("image_slots", 0),
            slide.get("has_charts", False),
            slide.get("has_tables", False),
            visual.get("layout_family"),
            visual.get("column_count"),
            visual.get("image_alignment"),
            visual.get("emphasis"),
            role_signature,
        )
        if fingerprint in seen_fingerprints:
            continue
        seen_fingerprints.add(fingerprint)
        selected.append(slide)
    return selected
